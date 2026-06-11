#!/usr/bin/env python3
"""Convert the extracted course decks from 4:3 (10 x 7.5 in) to 16:9
widescreen (13.333 x 7.5 in).

Heights and vertical positions are untouched. Horizontal positions and
widths are stretched by 4/3 so full-width header bars and text boxes keep
spanning the slide; pictures keep their aspect ratio (same size,
re-centred at their stretched centre point). Slide masters and layouts get
the same treatment so placeholder inheritance stays consistent.

Usage: python widescreen.py deck1.pptx [deck2.pptx ...]
"""
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

NEW_W = 12192000  # 13.333 in, EMU
OLD_W = 9144000   # 10 in


def stretch(shapes, factor):
    for sh in shapes:
        if sh.left is None or sh.width is None:
            continue
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            centre = sh.left + sh.width / 2
            sh.left = int(centre * factor - sh.width / 2)
        else:
            sh.left = int(sh.left * factor)
            sh.width = int(sh.width * factor)


def clamp(shapes):
    """Snap shapes that overhang the slide back onto the canvas.

    Some legacy .ppt decks carry text placeholders wider than the slide
    itself; PowerPoint masks the overhang, other renderers do not, and
    after stretching the centred text lands off-screen. Re-centre any
    such shape at full usable width.
    """
    margin = int(0.33 * 914400)
    for sh in shapes:
        if sh.left is None or sh.width is None:
            continue
        if sh.width > NEW_W or sh.left + sh.width > NEW_W + margin:
            sh.left = margin
            sh.width = NEW_W - 2 * margin
            print(f"  clamped overhanging shape {sh.name!r}")
        # zero-height text shapes render with the text vertically centred
        # on their top edge, clipping the upper half — give them room
        if (sh.has_text_frame and sh.text_frame.text.strip()
                and sh.height is not None and sh.height < int(0.3 * 914400)):
            sh.height = int(1.0 * 914400)
            if sh.top is not None and sh.top < int(0.1 * 914400):
                sh.top = int(0.15 * 914400)
            print(f"  grew degenerate-height text shape {sh.name!r}")


def convert(path):
    prs = Presentation(path)
    already = prs.slide_width == NEW_W
    factor = NEW_W / prs.slide_width
    if not already:
        prs.slide_width = NEW_W
    for master in prs.slide_masters:
        if not already:
            stretch(master.shapes, factor)
        clamp(master.shapes)
        for layout in master.slide_layouts:
            if not already:
                stretch(layout.shapes, factor)
            clamp(layout.shapes)
    for slide in prs.slides:
        if not already:
            stretch(slide.shapes, factor)
        clamp(slide.shapes)
    prs.save(path)
    print(f"{path}: {'re-clamped' if already else 'converted'} "
          f"to 13.333 x {prs.slide_height/914400:.2f} in")


for p in sys.argv[1:]:
    convert(p)
