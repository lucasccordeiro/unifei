#!/usr/bin/env python3
"""One-shot: append authored workshop appendices to the two extracted
Session 2 reference decks —

  session2/02-memory-model.pptx   4 -> 9 slides (Lab 2 background)
  session2/03-k-induction.pptx    5 -> 10 slides (Lab 3 stretch)

The base decks were extracted from COMP63342 and cannot be regenerated
here, so this edits them in place. The guard below refuses to run on
decks that are not at the pristine extracted slide counts, so re-running
is a no-op error, not duplicate slides. If the decks are ever
re-extracted, run this once afterwards.

Usage:  venv/bin/python extend_backup_decks.py
"""
import sys

from pptx import Presentation

from pptx_helpers import add_slide, add_body, add_table


DECKS = {"session2/02-memory-model.pptx": 4,
         "session2/03-k-induction.pptx": 5}


def open_pristine(path):
    """Open a deck, refusing to touch one that is not a pristine extract."""
    prs = Presentation(path)
    if len(prs.slides) != DECKS[path]:
        sys.exit(f"{path}: expected {DECKS[path]} slides (pristine "
                 f"extract), found {len(prs.slides)} — appendix already "
                 f"applied? Aborting.")
    if prs.slide_layouts[6].name != "Blank":
        sys.exit(f"{path}: slide layout 6 is "
                 f"{prs.slide_layouts[6].name!r}, not 'Blank' — "
                 f"pptx_helpers.add_slide would misrender. Aborting.")
    return prs


def extend_memory_model(prs):
    """Append the 5-slide Lab 2 appendix (slides 5-9)."""
    path = "session2/02-memory-model.pptx"

    s = add_slide(prs, "From this model to Lab 2")
    add_body(s, [
        "The slides before this one are the theory; tonight's files are "
        "the practice:",
        ("bounds encoding  →  overflow.c, vla.c", 1),
        ("allocation tracking  →  leak.c", 1),
        ("library models (gets)  →  getpassword.c", 1),
        "",
        "The next four slides say what ESBMC checks automatically, in "
        "workshop words.",
    ])

    s = add_slide(prs, "Every access gets a guard")
    add_body(s, [
        ("int a[2];  ...  a[i] = 0;", 0, True),
        "",
        "For every dereference ESBMC adds checks like:",
        ("0 <= i && i < 2          /* index inside the object */",
         1, True),
        ("p points into a live object", 1),
        "",
        "A failed guard prints as:",
        ("dereference failure: array bounds violated", 1, True),
        "",
        "Uninitialised i means the solver may pick ANY of the 2³² "
        "values — including -1 (overflow.c).",
    ])

    s = add_slide(prs, "What --memory-leak-check adds")
    add_body(s, [
        "Each malloc creates a tracked dynamic object.",
        "At every program exit: each still-allocated object must still "
        "be reachable through some pointer.",
        "",
        "When no pointer to a live object remains, the trace names it:",
        ("dereference failure: forgotten memory: dynamic_1_array",
         1, True),
        ("CWE: CWE-401", 1, True),
        ("Which statement in leak.c orphans that object? That is "
         "task 1.", 1),
        "",
        "A leak is unreachable memory — not merely un-freed memory.",
    ])

    s = add_slide(prs, "VLAs and symbolic sizes")
    add_body(s, [
        ("int a[n];   /* n is a runtime value — maybe nondet */",
         0, True),
        "",
        "Bounds checks become symbolic comparisons against n — the "
        "solver searches all sizes at once.",
        ("A zero-sized VLA is undefined behaviour before any element "
         "is touched. Does any call here create one?", 1),
        ("Your unwind bound must cover the loop that does the access "
         "— count its iterations before you pick.", 1),
    ])

    s = add_slide(prs, "The CWE labels, decoded")
    add_table(s, [
        ["CWE", "Meaning"],
        ["CWE-121", "Stack-based buffer overflow"],
        ["CWE-125", "Out-of-bounds read"],
        ["CWE-129", "Improper validation of array index"],
        ["CWE-131", "Incorrect calculation of buffer size"],
        ["CWE-193", "Off-by-one error"],
        ["CWE-787", "Out-of-bounds write"],
        ["CWE-401", "Missing release of memory after effective "
                     "lifetime (a memory leak)"],
    ], top=1.6, size=16, col_widths=[2.0, 8.5])
    add_body(s, ["One violated guard can imply several classes — ESBMC "
                 "prints them all (getpassword.c prints the first six)."],
             top=5.9, size=18)

    prs.save(path)
    print(f"saved {path} with {len(prs.slides)} slides")


def extend_k_induction(prs):
    """Append the 5-slide Lab 3 stretch appendix (slides 6-10)."""
    path = "session2/03-k-induction.pptx"

    s = add_slide(prs, "What --unwind k actually claims")
    add_body(s, [
        "SUCCESSFUL at bound k = no violation in any execution that "
        "fits within k loop unwindings.",
        "",
        "The unwinding assertion is the honesty check:",
        ("if it FAILS, the bound truncated the loop — a bound problem, "
         "not a program bug.", 1),
        ("if it holds, every execution was fully explored: the bounded "
         "verdict is a complete one.", 1),
    ])

    s = add_slide(prs, "unwind.c — the experiment")
    add_table(s, [
        ["Run", "Verdict", "What it means"],
        ["--unwind 5", "FAILED (unwinding assertion)",
         "Loop can run more than 5 times"],
        ["--unwind 20", "FAILED (unwinding assertion)", "Still short"],
        ["--unwind 60", "SUCCESSFUL",
         "Complete proof — every n <= MAX=50 fits"],
        ["--k-induction", "UNKNOWN", "See the next two slides"],
    ], top=1.7, size=16, col_widths=[2.4, 3.9, 5.4])
    add_body(s, ["60 > MAX means no execution is truncated: the bounded "
                 "check quantifies over ALL inputs."], top=4.6, size=20)

    s = add_slide(prs, "k-induction: the three checks")
    add_body(s, [
        "Prove a loop safe for ANY bound with three bounded queries:",
        ("Base case — no violation in the first k iterations.", 1),
        ("Forward condition — k iterations reach every reachable "
         "state (then the base case alone is a proof).", 1),
        ("Inductive step — from ANY k safe consecutive iterations, "
         "iteration k+1 is also safe.", 1),
        "",
        "Each check is a bounded (finite) formula — same machinery as "
        "BMC, different question.",
    ])

    s = add_slide(prs, "Why UNKNOWN on unwind.c")
    add_body(s, [
        "The inductive step starts from a havocked state: i and sum are "
        "ARBITRARY values, not values the program can reach.",
        ("From a mid-loop state where sum != 2*i, the step cannot "
         "conclude sum == 2*n at exit.", 1),
        ("So the step fails — for every k. ESBMC answers:", 1),
        ("VERIFICATION UNKNOWN", 1, True),
        "",
        "The assertion alone is not k-inductive. Nothing is wrong with "
        "the program — the proof needs help.",
    ])

    s = add_slide(prs, "The missing invariant")
    add_body(s, [
        ("sum == 2 * i      /* holds at every loop head */", 0, True),
        "",
        "Add it and the inductive step closes: a havocked state must "
        "now satisfy the invariant, and the exit assertion follows "
        "(i == n there).",
        "",
        "Stating loop invariants is exactly what deductive verifiers "
        "ask of you:",
        ("Dafny, Frama-C/ACSL — invariant annotations checked by a "
         "solver.", 1),
        ("BMC needed no annotations but gave a bounded answer; "
         "invariants buy the unbounded one.", 1),
    ])

    prs.save(path)
    print(f"saved {path} with {len(prs.slides)} slides")


if __name__ == "__main__":
    # Open (and guard) both decks BEFORE mutating either, so a mixed
    # state — one extended, one pristine — aborts cleanly instead of
    # leaving the pristine one forever unrepairable.
    decks = {path: open_pristine(path) for path in DECKS}
    extend_memory_model(decks["session2/02-memory-model.pptx"])
    extend_k_induction(decks["session2/03-k-induction.pptx"])
