#!/usr/bin/env python3
"""Build workshop-extras.pptx — the authored slides that complement the
three decks extracted from the COMP63342 lectures (see slides/README.md).

The canonical content lives in workshop-extras.md; this script renders the
same content as a native, editable PowerPoint file.

Usage:  python3 -m venv venv && venv/bin/pip install python-pptx
        venv/bin/python build_extras.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

PURPLE = RGBColor(0x5A, 0x0A, 0x7A)
DARK = RGBColor(0x21, 0x21, 0x21)
MONO = "Courier New"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide(title, title_size=36, lead=False):
    s = prs.slides.add_slide(BLANK)
    top = Inches(2.6) if lead else Inches(0.4)
    box = s.shapes.add_textbox(Inches(0.6), top, Inches(12.1), Inches(1.1))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(44 if lead else title_size)
    r.font.bold = True
    r.font.color.rgb = PURPLE
    return s


def add_body(slide, items, *, top=1.6, size=22, left=0.8, width=11.8):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(7.5 - top - 0.3))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        level = 0
        mono = False
        text = item
        if isinstance(item, tuple):
            text = item[0]
            level = item[1] if len(item) > 1 else 0
            mono = len(item) > 2 and item[2]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(0 if mono else 8)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size - 2 * level)
        r.font.color.rgb = DARK
        if mono:
            r.font.name = MONO
            r.font.size = Pt(size - 4)
    return box


def add_table(slide, rows, *, top=1.7, size=18, col_widths=None, left=0.8):
    n_rows, n_cols = len(rows), len(rows[0])
    total_w = sum(col_widths) if col_widths else 11.7
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top),
                                   Inches(total_w), Inches(0.4 * n_rows))
    table = shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for i, row in enumerate(rows):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = cell_text
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(size)
                    if i == 0:
                        r.font.bold = True
    return table


# ---------------------------------------------------------------- Session 1
s = add_slide("Formal Verification Workshop", lead=True)
add_body(s, ["Session 1 — Fundamentals and Real-World Applications",
             "18:00 – 21:00",
             "Built on the COMP63342 Software Security course "
             "(University of Manchester)"], top=3.9, size=26)

s = add_slide("Tonight")
add_table(s, [
    ["18:00", "Welcome + a question for you"],
    ["18:10", "Why software fails"],
    ["18:35", "Exercise 1: spot the bug"],
    ["18:50", "Testing vs. verification"],
    ["19:15", "Break"],
    ["19:25", "From programs to formulas: SAT, SMT"],
    ["19:50", "Bounded model checking"],
    ["20:15", "Exercise 2: predict the verdict"],
    ["20:30", "Who uses this in the real world?"],
    ["20:50", "Quiz + what to install for Session 2"],
], top=1.5, size=20, col_widths=[1.6, 10.1])

s = add_slide("Icebreaker")
add_body(s, ["Name a software failure that made the news.",
             "Any failure. Any decade. Shout it out / drop it in the poll."],
         top=2.6, size=30)

s = add_slide("Exercise 1 — Spot the Bug (15 min)")
add_body(s, [
    "In pairs, 7 minutes, on the handout:",
    ("int *zPtr;", 0, True),
    ("int *aPtr = NULL;", 0, True),
    ("void *sPtr = NULL;", 0, True),
    ("int number, i;", 0, True),
    ("int z[5] = {1, 2, 3, 4, 5};", 0, True),
    ("sPtr = z;   ++zPtr;   number = zPtr;", 0, True),
    ("number = *zPtr[2];   number = *sPtr;   ++z;", 0, True),
    "For each bug: WHAT is wrong, and what could happen at runtime?",
    "There are at least six. Debrief: one bug per pair.",
])

s = add_slide("Where does your technique sit?")
add_table(s, [
    ["", "Complete (no false alarms)", "Incomplete"],
    ["Sound (no missed bugs)", "the dream", "abstract interpretation"],
    ["Unsound", "bounded model checking*", "testing, code review"],
], top=1.7, size=20, col_widths=[3.4, 4.4, 3.9])
add_body(s, ["* sound up to its bound — tonight's central fine print",
             "Think–pair–share (5 min): place testing, code review, compiler "
             "warnings, and “prove it with maths” on this grid. "
             "Defend your placement."], top=4.1, size=22)

s = add_slide("Exercise 2 — Predict the Verdict (15 min)")
add_body(s, [
    "Groups of 3–4. Three short C programs on the handout "
    "(predict1.c, predict2.c, predict3.c).",
    "For each program, write down BEFORE we run anything:",
    ("Will the model checker say VERIFICATION FAILED or SUCCESSFUL?", 1),
    ("If FAILED — which line is the culprit?", 1),
    "Then we run ESBMC live and settle every bet.",
])

s = add_slide("Who uses this? — Model checking")
add_body(s, [
    "Amazon Web Services — CBMC proves memory safety of the TLS library "
    "s2n; bounded proofs run in CI on every commit",
    "ESBMC / CBMC — the tools from tonight's demos; compete yearly in "
    "SV-COMP on tens of thousands of benchmarks",
    "The counterexample you saw for getpassword.c is the same artefact "
    "AWS engineers read when a proof fails",
])

s = add_slide("Who uses this? — Symbolic execution & fuzzing")
add_body(s, [
    "KLEE found 56 bugs in GNU coreutils — code tested for 15 years",
    "Microsoft SAGE fuzzed Windows file parsers with symbolic execution; "
    "credited with finding 1/3 of all Win7 file-parsing bugs",
    "AFL / OSS-Fuzz — Google fuzzes ~1,000 open-source projects "
    "continuously; tens of thousands of bugs found",
    "The pragmatic cousin of verification: no proofs, ruthless "
    "effectiveness",
])

s = add_slide("Who uses this? — Proofs all the way up")
add_body(s, [
    "Astrée (abstract interpretation) proved absence of runtime errors "
    "in Airbus A380 fly-by-wire control code",
    "seL4 — an OS microkernel with a machine-checked proof of functional "
    "correctness",
    "CompCert — a C compiler with a proof that compilation preserves "
    "semantics; fuzzers find zero wrong-code bugs in it",
])

s = add_slide("Discussion")
add_body(s, [
    "If verification can prove the absence of bugs…",
    "why isn't all software verified?",
    "(Let's collect reasons — then I'll tell you which ones the industry "
    "actually cites.)",
], top=2.4, size=28)

s = add_slide("Closing quiz")
add_body(s, [
    "1. A tool that never raises a false alarm is called …?",
    "2. ESBMC returns UNSAT for C ∧ ¬P at --unwind 10. "
    "What do we know?",
    "3. What's inside a counterexample?",
    "4. Which technique proved the A380 flight-control code free of "
    "runtime errors?",
    "5. Why can't testing prove the getPassword program safe?",
], size=24)

s = add_slide("Before Session 2 — 10 minutes of homework")
add_body(s, [
    "Follow handouts/setup.md in the repo:",
    ("1.  pip install z3-solver", 0, True),
    ("2.  download ESBMC: github.com/esbmc/esbmc/releases", 0, True),
    ("3.  smoke test:  esbmc labs/lab4/float.c", 0, True),
    ("    → should print VERIFICATION FAILED "
     "(yes, failed — that's the point)", 0, True),
    "Can't make it work? Come at 17:45 — we'll fix it together.",
])

# ---------------------------------------------------------------- Session 2
s = add_slide("Formal Verification Workshop", lead=True)
add_body(s, ["Session 2 — Hands-On", "18:00 – 21:00",
             "Tonight you drive."], top=3.9, size=26)

s = add_slide("Recap quiz (while machines boot)")
add_body(s, [
    "1. To prove property P with a solver, you ask whether … is "
    "satisfiable?",
    "2. sat + a model in that setup means …?",
    "3. --unwind 5 reports a violated unwinding assertion. "
    "Is the program buggy?",
    "",
    "Smoke test, everyone:",
    ("esbmc lab4/float.c   →  VERIFICATION FAILED", 0, True),
], size=24)

s = add_slide("Tonight")
add_table(s, [
    ["18:15", "Lab 1", "Constraint solving with Z3"],
    ["18:50", "Lab 2", "Bug hunting with ESBMC"],
    ["19:30", "Break", ""],
    ["19:40", "Lab 3", "Writing your own specifications"],
    ["20:15", "Lab 4", "Team challenge"],
    ["20:45", "", "Debrief + where next"],
], top=1.6, size=20, col_widths=[1.6, 1.8, 8.3])
add_body(s, ["Work in pairs · keep labs/CHEATSHEET.md open · "
             "predict before you run"], top=4.6, size=22)

s = add_slide("Lab 1 — Constraint solving (35 min)")
add_body(s, [
    "labs/lab1/lab1.py — three stages, instructions inside.",
    "The one trick behind the whole evening:",
    ("to prove P, ask the solver if ¬P is satisfiable.", 1),
    ("unsat ⇒ P holds.   A model ⇒ a counterexample.", 1),
    "Checkpoint before Lab 2: what does unsat mean in Stage 1 vs. "
    "Stage 3?",
])

s = add_slide("Lab 2 — Bug hunting (40 min)")
add_body(s, ["predict → verify → READ THE COUNTEREXAMPLE → "
             "fix → re-verify"], top=1.5, size=22)
add_table(s, [
    ["#", "File", "By"],
    ["1", "overflow.c", "—"],
    ["2", "leak.c", "19:05"],
    ["3", "vla.c  (two bugs!)", "19:20"],
    ["4", "getpassword.c", "stretch"],
], top=2.3, size=18, col_widths=[0.8, 5.2, 2.2])
add_body(s, ["Done = VERIFICATION SUCCESSFUL on your fixed file, same "
             "flags.",
             "Finished early? You are now a TA — help the pair next to "
             "you."], top=5.2, size=20)

s = add_slide("Lab 3 — Specifications (35 min)")
add_body(s, [
    "triangle.c verifies SUCCESSFUL as shipped.",
    "Does that mean it's correct?",
    "Write the two TODO properties and find out.",
    "",
    "The tool checks what you specify — with a weak spec, green is "
    "cheap.",
], top=2.0, size=28)

s = add_slide("Lab 3, Part 2 — what did we prove?")
add_body(s, [
    ("esbmc unwind.c --unwind 5      →  ?", 0, True),
    ("esbmc unwind.c --unwind 20     →  ?", 0, True),
    ("esbmc unwind.c --unwind 60     →  ?", 0, True),
    "",
    "What is the unwinding assertion telling you at 5 and 20?",
    "At 60 — what exactly is proven, and for which n?",
    "Stretch: --k-induction answers VERIFICATION UNKNOWN. Why is the "
    "unbounded question fundamentally harder?",
])

s = add_slide("Lab 4 — Team challenge (30 min)")
add_body(s, [
    "Teams of 3–4. Points on the whiteboard:",
    ("1 pt — verdict predicted correctly in writing, before running", 1),
    ("2 pts — file fixed and re-verified", 1),
    "",
    "Track A:  race.c — write down the failing interleaving first, then",
    ("esbmc race.c --context-bound 2", 0, True),
    "Track B:  float.c — primary-school arithmetic. Or is it?",
])

s = add_slide("Debrief")
add_body(s, [
    "One thing the tool caught that YOU wouldn't have?",
    "The arc you just walked: intuition → concepts → tools "
    "→ application",
    "",
    "Where next:",
    ("SV-COMP benchmarks: sv-comp.sosy-lab.org", 1),
    ("Full course: COMP63342 Software Security (slides online)", 1),
    ("ESBMC: esbmc.org · Z3: github.com/Z3Prover/z3", 1),
    ("The complementary technique: fuzzing (AFL, OSS-Fuzz)", 1),
    "",
    "Feedback poll: pace / best lab / one improvement. Thank you!",
])

prs.save("workshop-extras.pptx")
print(f"saved workshop-extras.pptx with {len(prs.slides)} slides")
