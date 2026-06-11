"""Shared helpers for building the authored workshop decks
(build_extras.py, build_session2.py) with python-pptx, in 16:9."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

PURPLE = RGBColor(0x5A, 0x0A, 0x7A)
DARK = RGBColor(0x21, 0x21, 0x21)
MONO = "Courier New"


def new_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_slide(prs, title, *, title_size=36, lead=False):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    top = Inches(2.6) if lead else Inches(0.4)
    box = s.shapes.add_textbox(Inches(0.6), top, Inches(12.1), Inches(1.1))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(44 if lead else title_size)
    r.font.bold = True
    r.font.color.rgb = PURPLE
    return s


def add_body(slide, items, *, top=1.6, size=22, left=0.8, width=11.8):
    """items: strings, or tuples (text, level) / (text, level, mono)."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(7.5 - top - 0.3))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        level, mono, text = 0, False, item
        if isinstance(item, tuple):
            text = item[0]
            level = item[1] if len(item) > 1 else 0
            mono = len(item) > 2 and item[2]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(0 if mono else 8)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size - 2 * level)
        r.font.color.rgb = DARK
        if mono:
            r.font.name = MONO
            r.font.size = Pt(size - 4)
    return box


def add_table(slide, rows, *, top=1.7, size=18, col_widths=None, left=0.8):
    n_rows, n_cols = len(rows), len(rows[0])
    total_w = sum(col_widths) if col_widths else 11.7
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top),
                                   Inches(total_w), Inches(0.4 * n_rows))
    table = shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for i, row in enumerate(rows):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = cell_text
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(size)
                    if i == 0:
                        r.font.bold = True
    return table
